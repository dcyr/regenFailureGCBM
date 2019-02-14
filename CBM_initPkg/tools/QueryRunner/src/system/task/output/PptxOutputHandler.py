# core
import os

# contrib
from pptx import Presentation
from pptx.text.text import MSO_AUTO_SIZE
from pptx.util import Pt

# QueryRunner
from AbstractOutputHandler import AbstractOutputHandler

class PptxOutputHandler(AbstractOutputHandler):
    '''
    Writes query output to a PowerPoint file.
    
    :param outputPath: the output file to create
    :type outputPath: str
    '''
    
    def __init__(self, outputPath, slideNum, shapeName):
        if not os.path.isfile(outputPath):
            raise IOError("File not found: {0}".format(outputPath))
        
        self.__outputPath = outputPath
        self.__slideNum = slideNum
        self.__shapeName = shapeName


    def __getTargetShape(self, slide):
        for shape in slide.shapes:
            if shape.name == self.__shapeName:
                return shape

        raise Exception("Shape '{0}' not found in slide {1}"
                        .format(self.__shapeName, self.__slideNum))


    def handle(self, dbTitle, output):
        '''
        See :meth:`.AbstractOutputHandler.handle`.
        '''
        outString = "\n".join([" | ".join("{0}".format(value) for value in row)
                               for row in output])
        
        ppt = Presentation(self.__outputPath)
        slideIndex = self.__slideNum - 1
        slide = ppt.slides[slideIndex]
        shape = self.__getTargetShape(slide)
        shape.text = outString
        textFrame = shape.text_frame
        textFrame.paragraphs[0].font.size = Pt(12)
        textFrame.word_wrap = True
        textFrame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        ppt.save(self.__outputPath)
